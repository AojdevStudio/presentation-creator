"""
Tests for the chart API module.
"""
import pytest
import pandas as pd
import numpy as np
from pathlib import Path
from pptx import Presentation
from src.visualization.chart_api import ChartAPI

@pytest.fixture
def sample_data():
    """Create sample data for testing."""
    return pd.DataFrame({
        'Category': ['A', 'B', 'C', 'D', 'E'],
        'Values': [10, 20, 15, 25, 30],
        'Series2': [5, 15, 10, 20, 25],
        'Group': ['X', 'X', 'Y', 'Y', 'Z']
    }).set_index('Category')

@pytest.fixture
def chart_api():
    """Create a test chart API instance."""
    return ChartAPI()

def test_select_data_range_basic(chart_api, sample_data):
    """Test basic data range selection."""
    result = chart_api.select_data_range(
        sample_data,
        start=1,
        end=3
    )
    assert len(result) == 2
    assert list(result.index) == ['B', 'C']

def test_select_data_range_columns(chart_api, sample_data):
    """Test data range selection with column filtering."""
    result = chart_api.select_data_range(
        sample_data,
        columns=['Values']
    )
    assert list(result.columns) == ['Values']
    assert len(result) == len(sample_data)

def test_select_data_range_conditions(chart_api, sample_data):
    """Test data range selection with conditions."""
    result = chart_api.select_data_range(
        sample_data,
        conditions={
            'Group': 'X'
        }
    )
    assert len(result) == 2
    assert all(idx in ['A', 'B'] for idx in result.index)

def test_select_data_range_complex_conditions(chart_api, sample_data):
    """Test data range selection with complex conditions."""
    result = chart_api.select_data_range(
        sample_data,
        conditions={
            'Values': {'operator': '>', 'value': 15},
            'Group': ['X', 'Y']
        }
    )
    assert len(result) == 2
    assert all(result['Values'] > 15)

def test_create_chart_basic(chart_api):
    """Test basic chart creation."""
    slide = chart_api.presentation.slides.add_slide(
        chart_api.presentation.slide_layouts[5]
    )
    
    data = {'A': [1, 2, 3], 'B': [4, 5, 6]}
    shape = chart_api.create_chart(
        slide,
        data,
        'bar',
        (1, 1, 6, 4),
        title="Test Chart"
    )
    
    assert shape is not None
    assert hasattr(shape, 'chart')  # Should be a native chart

def test_create_chart_with_data_selection(chart_api, sample_data):
    """Test chart creation with data selection."""
    slide = chart_api.presentation.slides.add_slide(
        chart_api.presentation.slide_layouts[5]
    )
    
    shape = chart_api.create_chart(
        slide,
        sample_data,
        'bar',
        (1, 1, 6, 4),
        title="Test Chart",
        data_selection={
            'columns': ['Values'],
            'conditions': {'Group': 'X'}
        }
    )
    
    assert shape is not None
    assert hasattr(shape, 'chart')

def test_quick_chart(chart_api, sample_data):
    """Test quick chart creation."""
    slide = chart_api.quick_chart(
        sample_data,
        'bar',
        title="Quick Chart",
        data_selection={'columns': ['Values']}
    )
    
    assert slide is not None
    assert len(slide.shapes) > 0
    # Should have at least one shape (the chart)
    chart_shapes = [shape for shape in slide.shapes if hasattr(shape, 'chart')]
    assert len(chart_shapes) == 1

def test_create_dashboard(chart_api, sample_data):
    """Test dashboard creation with multiple charts."""
    charts = [
        {
            'type': 'bar',
            'position': (1, 1, 4, 3),
            'title': 'Bar Chart',
            'data_selection': {'columns': ['Values']}
        },
        {
            'type': 'line',
            'position': (5, 1, 4, 3),
            'title': 'Line Chart',
            'data_selection': {'columns': ['Series2']}
        },
        {
            'type': 'pie',
            'position': (1, 4, 8, 3),
            'title': 'Pie Chart',
            'data_selection': {
                'columns': ['Values'],
                'conditions': {'Group': 'X'}
            }
        }
    ]
    
    slide = chart_api.create_dashboard(
        sample_data,
        charts,
        title="Test Dashboard"
    )
    
    assert slide is not None
    assert len(slide.shapes) >= len(charts)  # Should have at least as many shapes as charts

def test_save_presentation(chart_api, sample_data, tmp_path):
    """Test saving the presentation."""
    # Create a chart
    chart_api.quick_chart(
        sample_data,
        'bar',
        title="Test Chart"
    )
    
    # Save the presentation
    output_path = tmp_path / "test_output.pptx"
    saved_path = chart_api.save_presentation(output_path)
    
    assert Path(saved_path).exists()
    assert Path(saved_path).stat().st_size > 0

def test_error_handling(chart_api):
    """Test error handling for invalid inputs."""
    slide = chart_api.presentation.slides.add_slide(
        chart_api.presentation.slide_layouts[5]
    )
    
    # Test with empty data
    with pytest.raises(ValueError):
        chart_api.create_chart(
            slide,
            {},
            'bar',
            (1, 1, 6, 4)
        )
    
    # Test with invalid chart type
    with pytest.raises(ValueError):
        chart_api.create_chart(
            slide,
            {'A': [1, 2, 3]},
            'invalid_type',
            (1, 1, 6, 4)
        )
        
    # Test with invalid data selection
    with pytest.raises(KeyError):
        chart_api.create_chart(
            slide,
            {'A': [1, 2, 3]},
            'bar',
            (1, 1, 6, 4),
            data_selection={'columns': ['NonExistentColumn']}
        ) 