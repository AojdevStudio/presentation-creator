"""
Tests for the presentation finalizer module.
"""
import pytest
from datetime import datetime
from pptx import Presentation
from pptx.presentation import Presentation as PresentationType
from pptx.dml.color import RGBColor
from src.presentation.presentation_finalizer import PresentationFinalizer

@pytest.fixture
def presentation() -> PresentationType:
    """Create a test presentation."""
    return Presentation()

@pytest.fixture
def finalizer(presentation) -> PresentationFinalizer:
    """Create a test finalizer."""
    return PresentationFinalizer(presentation)

def test_add_metadata(finalizer):
    """Test adding metadata to presentation."""
    metadata = {
        'title': 'Test Presentation',
        'author': 'Test Author',
        'company': 'Test Company',
        'keywords': 'test, presentation',
        'comments': 'Test comments'
    }
    
    finalizer.add_metadata(metadata)
    
    core_props = finalizer.presentation.core_properties
    assert core_props.title == 'Test Presentation'
    assert core_props.author == 'Test Author'
    assert core_props.company == 'Test Company'
    assert core_props.keywords == 'test, presentation'
    assert core_props.comments == 'Test comments'
    assert isinstance(core_props.created, datetime)
    assert isinstance(core_props.modified, datetime)
    assert core_props.last_modified_by == 'Test Author'

def test_add_title_slide(finalizer):
    """Test adding a title slide."""
    finalizer.add_title_slide('Test Title', 'Test Subtitle')
    
    # Check if slide was added
    assert len(finalizer.presentation.slides) == 1
    
    slide = finalizer.presentation.slides[0]
    
    # Check title
    assert slide.shapes.title.text == 'Test Title'
    title_font = slide.shapes.title.text_frame.paragraphs[0].font
    assert title_font.size.pt == 44
    assert title_font.bold
    assert title_font.color.rgb == RGBColor(33, 33, 33)
    
    # Check subtitle
    subtitle_shape = slide.placeholders[1]
    assert subtitle_shape.text == 'Test Subtitle'
    subtitle_font = subtitle_shape.text_frame.paragraphs[0].font
    assert subtitle_font.size.pt == 24
    assert subtitle_font.color.rgb == RGBColor(89, 89, 89)

def test_add_footer(finalizer):
    """Test adding footer to slides."""
    # Add a few slides
    for _ in range(3):
        layout = finalizer.presentation.slide_layouts[0]
        finalizer.presentation.slides.add_slide(layout)
    
    footer_text = 'Test Footer'
    finalizer.add_footer(footer_text)
    
    # Check each slide has footer (except title slide)
    slides = list(finalizer.presentation.slides)
    for slide in slides[1:]:  # Skip title slide
        # Find footer shape by name
        footer_shapes = [shape for shape in slide.shapes 
                        if hasattr(shape, 'name') and shape.name == 'Footer Text']
        assert len(footer_shapes) == 1
        footer_shape = footer_shapes[0]
        assert footer_shape.text_frame.text == footer_text
        font = footer_shape.text_frame.paragraphs[0].font
        assert font.size.pt == 10
        assert font.color.rgb == RGBColor(128, 128, 128)

def test_add_slide_numbers(finalizer):
    """Test adding slide numbers."""
    # Add a few slides
    for _ in range(3):
        layout = finalizer.presentation.slide_layouts[0]
        finalizer.presentation.slides.add_slide(layout)
    
    finalizer.add_slide_numbers()
    
    # Skip first slide (title slide)
    slides = list(finalizer.presentation.slides)
    for i, slide in enumerate(slides[1:], start=2):
        # Get all shapes and find the slide number shape (should be a textbox)
        number_shapes = [shape for shape in slide.shapes if hasattr(shape, 'text_frame')]
        assert len(number_shapes) > 0
        number_shape = number_shapes[-1]  # Last added textbox should be slide number
        assert number_shape.text_frame.text == str(i)
        font = number_shape.text_frame.paragraphs[0].font
        assert font.size.pt == 10
        assert font.color.rgb == RGBColor(128, 128, 128)

def test_finalize(finalizer):
    """Test complete finalization process."""
    metadata = {
        'title': 'Test Presentation',
        'subtitle': 'Test Subtitle',
        'author': 'Test Author',
        'company': 'Test Company',
        'keywords': 'test, presentation',
        'comments': 'Test comments',
        'footer': 'Test Footer'
    }
    
    # Add some content slides
    for _ in range(2):
        layout = finalizer.presentation.slide_layouts[1]  # Content layout
        finalizer.presentation.slides.add_slide(layout)
    
    finalizer.finalize(metadata)
    
    # Check title slide was added first
    slides = list(finalizer.presentation.slides)
    assert len(slides) == 3  # Title slide + 2 content slides
    title_slide = slides[0]
    assert title_slide.shapes.title.text == 'Test Presentation'
    
    # Check metadata was added
    core_props = finalizer.presentation.core_properties
    assert core_props.title == 'Test Presentation'
    assert core_props.author == 'Test Author'
    
    # Check footer on all slides (except title slide)
    for slide in slides[1:]:
        # Find footer shape by name
        footer_shapes = [shape for shape in slide.shapes 
                        if hasattr(shape, 'name') and shape.name == 'Footer Text']
        assert len(footer_shapes) == 1
        footer_shape = footer_shapes[0]
        assert footer_shape.text_frame.text == 'Test Footer'
    
    # Check slide numbers (except title slide)
    for i, slide in enumerate(slides[1:], start=2):
        number_shapes = [shape for shape in slide.shapes if hasattr(shape, 'text_frame')]
        assert len(number_shapes) > 0
        number_shape = number_shapes[-1]  # Last shape should be slide number
        assert number_shape.text_frame.text == str(i) 