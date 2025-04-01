"""
Tests for the chart integrator module.
"""
import pytest
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import plotly.graph_objects as go
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from src.visualization.chart_integrator import ChartIntegrator
from src.visualization.chart_generator import ChartGenerator
from src.core.theme_manager import ThemeManager

@pytest.fixture
def theme_manager():
    """Create a test theme manager."""
    return ThemeManager({
        "colors": {
            "primary": "#1f77b4",
            "secondary": "#ff7f0e",
            "tertiary": "#2ca02c"
        },
        "fonts": {
            "title": "Arial",
            "body": "Calibri"
        },
        "spacing": {
            "paragraph": 1.15,
            "line": 1.0
        },
        "alignment": {
            "title": "center",
            "body": "left"
        }
    })

@pytest.fixture
def chart_generator(theme_manager):
    """Create a test chart generator."""
    return ChartGenerator(theme_manager)

@pytest.fixture
def presentation():
    """Create a test presentation."""
    return Presentation()

@pytest.fixture
def slide(presentation):
    """Create a test slide."""
    return presentation.slides.add_slide(presentation.slide_layouts[5])  # Blank layout

@pytest.fixture
def chart_integrator(presentation, chart_generator):
    """Create a test chart integrator."""
    return ChartIntegrator(presentation, chart_generator)

@pytest.fixture
def sample_data():
    """Create sample data for testing."""
    return pd.DataFrame({
        'Category': ['A', 'B', 'C', 'D'],
        'Values': [10, 20, 15, 25],
        'Series2': [5, 15, 10, 20]
    }).set_index('Category')

def test_add_native_chart(chart_integrator, slide, sample_data):
    """Test adding a native PowerPoint chart."""
    chart_shape = chart_integrator.add_native_chart(
        slide,
        sample_data,
        'bar',
        (1, 1, 6, 4),  # left, top, width, height in inches
        title="Test Bar Chart"
    )
    
    assert chart_shape is not None
    assert chart_shape.chart.chart_type == XL_CHART_TYPE.COLUMN_CLUSTERED
    assert chart_shape.chart.has_title
    assert chart_shape.chart.chart_title.text_frame.text == "Test Bar Chart"

def test_add_native_chart_invalid_type(chart_integrator, slide, sample_data):
    """Test adding a native chart with invalid type."""
    with pytest.raises(ValueError) as exc_info:
        chart_integrator.add_native_chart(
            slide,
            sample_data,
            'invalid_type',
            (1, 1, 6, 4)
        )
    assert "Unsupported native chart type" in str(exc_info.value)

def test_add_image_chart_matplotlib(chart_integrator, slide, sample_data):
    """Test adding a matplotlib chart as image."""
    # Create matplotlib figure
    fig, ax = plt.subplots()
    ax.plot(sample_data.index, sample_data['Values'])
    
    picture = chart_integrator.add_image_chart(
        slide,
        fig,
        (1, 1, 6, 4)
    )
    
    assert picture is not None
    plt.close(fig)

def test_add_image_chart_plotly(chart_integrator, slide, sample_data):
    """Test adding a plotly chart as image."""
    # Create plotly figure
    fig = go.Figure(data=go.Scatter(
        x=sample_data.index,
        y=sample_data['Values']
    ))
    
    picture = chart_integrator.add_image_chart(
        slide,
        fig,
        (1, 1, 6, 4)
    )
    
    assert picture is not None

def test_add_chart_native(chart_integrator, slide, sample_data):
    """Test adding a chart with native PowerPoint format."""
    shape = chart_integrator.add_chart(
        slide,
        sample_data,
        'bar',
        (1, 1, 6, 4),
        title="Test Chart",
        use_native=True
    )
    
    assert shape is not None
    assert shape.chart.chart_type == XL_CHART_TYPE.COLUMN_CLUSTERED

def test_add_chart_image(chart_integrator, slide, sample_data):
    """Test adding a chart as image."""
    shape = chart_integrator.add_chart(
        slide,
        sample_data,
        'combo',  # Complex chart type that requires image-based approach
        (1, 1, 6, 4),
        title="Test Chart",
        bar_columns=['Values'],
        line_columns=['Series2']
    )
    
    assert shape is not None

def test_add_chart_auto_selection(chart_integrator, slide, sample_data):
    """Test automatic selection between native and image-based charts."""
    # Should use native chart for basic bar chart
    shape1 = chart_integrator.add_chart(
        slide,
        sample_data,
        'bar',
        (1, 1, 6, 4)
    )
    assert hasattr(shape1, 'chart')
    
    # Should use image for combo chart
    shape2 = chart_integrator.add_chart(
        slide,
        sample_data,
        'combo',
        (1, 6, 6, 4),
        bar_columns=['Values'],
        line_columns=['Series2']
    )
    assert not hasattr(shape2, 'chart')

def test_chart_styling(chart_integrator, slide, sample_data):
    """Test chart styling options."""
    shape = chart_integrator.add_chart(
        slide,
        sample_data,
        'bar',
        (1, 1, 6, 4),
        title="Styled Chart",
        use_native=True,
        style=2,  # Apply chart style
        has_legend=True
    )
    
    assert shape.chart.has_legend
    assert shape.chart.style == 2

def test_error_handling(chart_integrator, slide):
    """Test error handling for invalid data."""
    with pytest.raises(ValueError):
        chart_integrator.add_chart(
            slide,
            {},  # Empty data
            'bar',
            (1, 1, 6, 4)
        ) 