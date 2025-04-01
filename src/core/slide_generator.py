"""
Slide generation functions using python-pptx.
"""
import logging
from typing import Dict, List, Optional, Any, Tuple
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from .theme_manager import ThemeManager

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class SlideGenerator:
    """Generates different types of slides using python-pptx"""
    
    def __init__(self, template_path: Optional[str] = None, custom_theme: Optional[Dict[str, Any]] = None):
        """Initialize the slide generator.
        
        Args:
            template_path: Optional path to a PowerPoint template file
            custom_theme: Optional custom theme settings
        """
        self.prs = Presentation(template_path) if template_path else Presentation()
        self.theme_manager = ThemeManager(custom_theme)
        
    def _add_text_box(self, slide, left: float, top: float, width: float, height: float,
                     text: str, style_type: str = "body") -> None:
        """Add a text box to a slide with themed styling.
        
        Args:
            slide: The slide to add the text box to
            left: Left position in inches
            top: Top position in inches
            width: Width in inches
            height: Height in inches
            text: Text content
            style_type: Type of text style to apply (title, subtitle, heading, body, caption)
        """
        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = txBox.text_frame
        tf.text = text
        self.theme_manager.apply_text_style(tf.paragraphs[0], style_type)
        
    def create_title_slide(self, content: Dict[str, Any]) -> None:
        """Create a title slide.
        
        Args:
            content: Dictionary containing title, subtitle, presenter, and date
        """
        slide_layout = self.prs.slide_layouts[0]  # Title slide layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Add title
        title = slide.shapes.title
        title.text = content["title"]
        self.theme_manager.apply_text_style(title.text_frame.paragraphs[0], "title")
        
        # Add subtitle
        subtitle = slide.placeholders[1]
        subtitle.text = content["subtitle"]
        self.theme_manager.apply_text_style(subtitle.text_frame.paragraphs[0], "subtitle")
        
        # Add presenter and date
        self._add_text_box(slide, 0.5, 5.0, 4.0, 0.5, 
                          f"Presenter: {content['presenter']}", 
                          style_type="caption")
        self._add_text_box(slide, 8.5, 5.0, 2.0, 0.5, 
                          content['date'], 
                          style_type="caption")
        
    def create_content_slide(self, content: Dict[str, Any]) -> None:
        """Create a content slide with bullet points.
        
        Args:
            content: Dictionary containing title and key points
        """
        slide_layout = self.prs.slide_layouts[1]  # Content slide layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Add title
        title = slide.shapes.title
        title.text = content["title"]
        self.theme_manager.apply_text_style(title.text_frame.paragraphs[0], "heading")
        
        # Add bullet points
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.text = ""  # Clear default text
        
        for point in content["key_points"]:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0  # Top level bullet
            self.theme_manager.apply_text_style(p, "body")
            
        # Add context if provided
        if "context" in content:
            self._add_text_box(slide, 0.5, 5.0, 9.0, 0.75,
                             content["context"],
                             style_type="caption")
            
    def create_section_transition(self, content: Dict[str, Any]) -> None:
        """Create a section transition slide.
        
        Args:
            content: Dictionary containing current and next section info
        """
        slide_layout = self.prs.slide_layouts[2]  # Section header layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Add current section
        title = slide.shapes.title
        title.text = content["current_section"]
        self.theme_manager.apply_text_style(title.text_frame.paragraphs[0], "title")
        
        # Add next section preview
        subtitle = slide.placeholders[1]
        subtitle.text = f"Coming Up: {content['next_section']}"
        self.theme_manager.apply_text_style(subtitle.text_frame.paragraphs[0], "subtitle")
        
    def create_summary_slide(self, content: Dict[str, Any]) -> None:
        """Create a summary slide.
        
        Args:
            content: Dictionary containing main topics and key takeaways
        """
        slide_layout = self.prs.slide_layouts[1]  # Content slide layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Add title
        title = slide.shapes.title
        title.text = "Summary"
        self.theme_manager.apply_text_style(title.text_frame.paragraphs[0], "heading")
        
        # Add main topics and takeaways
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.text = ""  # Clear default text
        
        # Add main topics
        p = tf.add_paragraph()
        p.text = "Main Topics:"
        self.theme_manager.apply_text_style(p, "heading")
        
        for topic in content["main_topics"]:
            p = tf.add_paragraph()
            p.text = topic
            p.level = 1
            self.theme_manager.apply_text_style(p, "body")
            
        # Add space between sections
        p = tf.add_paragraph()
        p.text = ""
            
        # Add key takeaways
        p = tf.add_paragraph()
        p.text = "Key Takeaways:"
        self.theme_manager.apply_text_style(p, "heading")
        
        for takeaway in content["key_takeaways"]:
            p = tf.add_paragraph()
            p.text = takeaway
            p.level = 1
            self.theme_manager.apply_text_style(p, "body")
            
    def save(self, output_path: str) -> None:
        """Save the presentation.
        
        Args:
            output_path: Path where to save the presentation
        """
        self.prs.save(output_path) 