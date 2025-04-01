"""
Module for integrating charts with PowerPoint presentations.
"""
from typing import Dict, List, Any, Optional, Union, Tuple
from pathlib import Path
import io
import matplotlib.pyplot as plt
import plotly.graph_objects as go
from pptx.presentation import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from .chart_generator import ChartGenerator

class ChartIntegrator:
    """Handles integration of charts into PowerPoint presentations."""
    
    def __init__(self, presentation: Presentation, chart_generator: Optional[ChartGenerator] = None):
        """
        Initialize the chart integrator.
        
        Args:
            presentation: PowerPoint presentation to add charts to
            chart_generator: Optional chart generator instance
        """
        self.presentation = presentation
        self.chart_generator = chart_generator or ChartGenerator()
        
    def add_native_chart(self,
                        slide: Any,
                        data: Union[Dict[str, List[Any]], List[Any]],
                        chart_type: str,
                        position: Tuple[float, float, float, float],
                        title: Optional[str] = None,
                        **kwargs) -> Any:
        """
        Add a native PowerPoint chart to a slide.
        
        Args:
            slide: PowerPoint slide to add chart to
            data: Chart data
            chart_type: Type of chart to create
            position: Chart position (left, top, width, height) in inches
            title: Optional chart title
            **kwargs: Additional chart options
            
        Returns:
            The created chart shape
        """
        # Map chart types to PowerPoint chart types
        chart_type_map = {
            'bar': XL_CHART_TYPE.COLUMN_CLUSTERED,
            'bar_stacked': XL_CHART_TYPE.COLUMN_STACKED,
            'bar_stacked_100': XL_CHART_TYPE.COLUMN_STACKED_100,
            'column': XL_CHART_TYPE.COLUMN_CLUSTERED,
            'column_stacked': XL_CHART_TYPE.COLUMN_STACKED,
            'column_stacked_100': XL_CHART_TYPE.COLUMN_STACKED_100,
            'line': XL_CHART_TYPE.LINE,
            'line_markers': XL_CHART_TYPE.LINE_MARKERS,
            'pie': XL_CHART_TYPE.PIE,
            'doughnut': XL_CHART_TYPE.DOUGHNUT,
            'radar': XL_CHART_TYPE.RADAR,
            'scatter': XL_CHART_TYPE.XY_SCATTER,
            'area': XL_CHART_TYPE.AREA,
            'area_stacked': XL_CHART_TYPE.AREA_STACKED,
            'area_stacked_100': XL_CHART_TYPE.AREA_STACKED_100
        }
        
        if chart_type not in chart_type_map:
            raise ValueError(f"Unsupported native chart type: {chart_type}")
            
        # Convert position to inches
        left, top, width, height = [Inches(x) for x in position]
        
        # Prepare chart data
        chart_data = CategoryChartData()
        df = self.chart_generator._prepare_data(data)
        
        # Add categories (index)
        categories = [str(x) for x in df.index]
        chart_data.categories = categories
        
        # Add series
        for column in df.columns:
            chart_data.add_series(column, list(df[column]))
            
        # Create chart
        chart_shape = slide.shapes.add_chart(
            chart_type_map[chart_type],
            left, top, width, height,
            chart_data
        )
        
        # Configure chart
        chart = chart_shape.chart
        if title:
            chart.has_title = True
            chart.chart_title.text_frame.text = title
            
        # Apply additional styling from kwargs
        if 'style' in kwargs:
            chart.style = kwargs['style']
        if 'has_legend' in kwargs:
            chart.has_legend = kwargs['has_legend']
            
        return chart_shape
        
    def add_image_chart(self,
                       slide: Any,
                       figure: Union[plt.Figure, go.Figure],
                       position: Tuple[float, float, float, float],
                       image_format: str = 'png') -> Any:
        """
        Add a chart as an image to a slide.
        
        Args:
            slide: PowerPoint slide to add chart to
            figure: Matplotlib or Plotly figure
            position: Image position (left, top, width, height) in inches
            image_format: Image format ('png' or 'svg')
            
        Returns:
            The created picture shape
        """
        # Convert figure to bytes
        image_data = self.chart_generator.save_figure_to_bytes(figure, format=image_format)
        
        # Convert position to inches
        left, top, width, height = [Inches(x) for x in position]
        
        # Create image stream
        image_stream = io.BytesIO(image_data)
        
        # Add picture to slide
        picture = slide.shapes.add_picture(
            image_stream,
            left, top, width, height
        )
        
        return picture
        
    def add_chart(self,
                  slide: Any,
                  data: Union[Dict[str, List[Any]], List[Any]],
                  chart_type: str,
                  position: Tuple[float, float, float, float],
                  title: Optional[str] = None,
                  use_native: Optional[bool] = None,
                  **kwargs) -> Any:
        """
        Add a chart to a slide, automatically choosing between native and image-based.
        
        Args:
            slide: PowerPoint slide to add chart to
            data: Chart data
            chart_type: Type of chart to create
            position: Chart position (left, top, width, height) in inches
            title: Optional chart title
            use_native: Force using native PowerPoint charts if True
            **kwargs: Additional chart options
            
        Returns:
            The created chart or picture shape
        """
        # Determine if we should use native charts
        if use_native is None:
            # Use native charts for basic chart types
            native_types = {'bar', 'column', 'line', 'pie', 'area', 'scatter'}
            use_native = chart_type.split('_')[0] in native_types
            
        if use_native:
            try:
                return self.add_native_chart(slide, data, chart_type, position, title, **kwargs)
            except ValueError as e:
                if "Unsupported native chart type" in str(e):
                    # Fall back to image-based chart
                    use_native = False
                else:
                    raise
                    
        if not use_native:
            # Create chart using chart generator
            if chart_type == 'scatter':
                fig = self.chart_generator.create_scatter_plot(
                    data, title=title,
                    x_column=kwargs.get('x_column', data.columns[0]),
                    y_column=kwargs.get('y_column', data.columns[1])
                )
            elif chart_type.startswith('bar'):
                fig = self.chart_generator.create_bar_chart(
                    data, title=title,
                    orientation='horizontal' if chart_type == 'bar' else 'vertical',
                    stacked='stacked' in chart_type
                )
            elif chart_type.startswith('line'):
                fig = self.chart_generator.create_line_chart(
                    data, title=title,
                    markers='markers' in chart_type
                )
            elif chart_type == 'pie':
                fig = self.chart_generator.create_pie_chart(
                    data, title=title
                )
            elif chart_type == 'combo':
                fig = self.chart_generator.create_combo_chart(
                    data, title=title,
                    bar_columns=kwargs.get('bar_columns', []),
                    line_columns=kwargs.get('line_columns', [])
                )
            else:
                raise ValueError(f"Unsupported chart type: {chart_type}")
                
            return self.add_image_chart(slide, fig, position) 