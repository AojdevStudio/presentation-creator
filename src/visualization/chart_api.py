"""
High-level API for chart generation and data range selection.
"""
from typing import Dict, List, Any, Optional, Union, Tuple
import pandas as pd
from pathlib import Path
from pptx.presentation import Presentation
from pptx.slide import Slide
from .chart_generator import ChartGenerator
from .chart_integrator import ChartIntegrator

class ChartAPI:
    """High-level API for chart generation and PowerPoint integration."""
    
    def __init__(self, presentation: Optional[Presentation] = None):
        """
        Initialize the chart API.
        
        Args:
            presentation: Optional PowerPoint presentation to work with
        """
        self.presentation = presentation or Presentation()
        self.chart_generator = ChartGenerator()
        self.chart_integrator = ChartIntegrator(self.presentation, self.chart_generator)
        
    def select_data_range(self,
                         data: Union[pd.DataFrame, Dict[str, List[Any]], List[Any]],
                         start: Optional[Union[int, str]] = None,
                         end: Optional[Union[int, str]] = None,
                         columns: Optional[List[str]] = None,
                         conditions: Optional[Dict[str, Any]] = None) -> pd.DataFrame:
        """
        Select a subset of data based on various criteria.
        
        Args:
            data: Input data
            start: Start index/row for selection
            end: End index/row for selection
            columns: List of columns to include
            conditions: Dictionary of filtering conditions
            
        Returns:
            Selected data as DataFrame
        """
        df = self.chart_generator._prepare_data(data)
        
        # Apply row selection
        if start is not None or end is not None:
            if isinstance(start, str) and start in df.index:
                start_idx = df.index.get_loc(start)
            else:
                start_idx = start or 0
                
            if isinstance(end, str) and end in df.index:
                end_idx = df.index.get_loc(end) + 1
            else:
                end_idx = end or len(df)
                
            df = df.iloc[start_idx:end_idx]
            
        # Apply column selection
        if columns:
            df = df[columns]
            
        # Apply filtering conditions
        if conditions:
            for column, condition in conditions.items():
                if isinstance(condition, (list, tuple)):
                    df = df[df[column].isin(condition)]
                elif isinstance(condition, dict):
                    operator = condition.get('operator', '==')
                    value = condition.get('value')
                    
                    if operator == '==':
                        df = df[df[column] == value]
                    elif operator == '!=':
                        df = df[df[column] != value]
                    elif operator == '>':
                        df = df[df[column] > value]
                    elif operator == '>=':
                        df = df[df[column] >= value]
                    elif operator == '<':
                        df = df[df[column] < value]
                    elif operator == '<=':
                        df = df[df[column] <= value]
                else:
                    df = df[df[column] == condition]
                    
        return df
        
    def create_chart(self,
                    slide: Slide,
                    data: Union[pd.DataFrame, Dict[str, List[Any]], List[Any]],
                    chart_type: str,
                    position: Tuple[float, float, float, float],
                    title: Optional[str] = None,
                    data_selection: Optional[Dict[str, Any]] = None,
                    chart_options: Optional[Dict[str, Any]] = None) -> Any:
        """
        Create and add a chart to a slide with optional data selection.
        
        Args:
            slide: PowerPoint slide to add chart to
            data: Input data
            chart_type: Type of chart to create
            position: Chart position (left, top, width, height) in inches
            title: Optional chart title
            data_selection: Optional data selection criteria
            chart_options: Optional chart configuration options
            
        Returns:
            Created chart or picture shape
        """
        # Select data if criteria provided
        if data_selection:
            data = self.select_data_range(
                data,
                start=data_selection.get('start'),
                end=data_selection.get('end'),
                columns=data_selection.get('columns'),
                conditions=data_selection.get('conditions')
            )
            
        # Add chart with options
        return self.chart_integrator.add_chart(
            slide,
            data,
            chart_type,
            position,
            title=title,
            **(chart_options or {})
        )
        
    def quick_chart(self,
                   data: Union[pd.DataFrame, Dict[str, List[Any]], List[Any]],
                   chart_type: str,
                   title: Optional[str] = None,
                   layout_index: int = 5,  # Default to blank layout
                   position: Optional[Tuple[float, float, float, float]] = None,
                   data_selection: Optional[Dict[str, Any]] = None,
                   chart_options: Optional[Dict[str, Any]] = None) -> Slide:
        """
        Quickly create a new slide with a chart.
        
        Args:
            data: Input data
            chart_type: Type of chart to create
            title: Optional chart title
            layout_index: Slide layout index to use
            position: Optional chart position (defaults to centered)
            data_selection: Optional data selection criteria
            chart_options: Optional chart configuration options
            
        Returns:
            Created slide
        """
        # Create new slide
        slide = self.presentation.slides.add_slide(
            self.presentation.slide_layouts[layout_index]
        )
        
        # Use default position if not specified
        if position is None:
            position = (1, 1, 8, 5)  # Reasonable default size and position
            
        # Create chart
        self.create_chart(
            slide,
            data,
            chart_type,
            position,
            title=title,
            data_selection=data_selection,
            chart_options=chart_options
        )
        
        return slide
        
    def create_dashboard(self,
                        data: Union[pd.DataFrame, Dict[str, List[Any]], List[Any]],
                        charts: List[Dict[str, Any]],
                        title: Optional[str] = None,
                        layout_index: int = 5) -> Slide:
        """
        Create a slide with multiple charts arranged as a dashboard.
        
        Args:
            data: Input data
            charts: List of chart specifications with positions and types
            title: Optional dashboard title
            layout_index: Slide layout index to use
            
        Returns:
            Created slide
        """
        # Create new slide
        slide = self.presentation.slides.add_slide(
            self.presentation.slide_layouts[layout_index]
        )
        
        # Add title if provided
        if title and hasattr(slide, 'shapes') and hasattr(slide.shapes, 'title'):
            title_shape = slide.shapes.title
            title_shape.text = title
            
        # Create each chart
        for chart_spec in charts:
            self.create_chart(
                slide,
                data,
                chart_spec['type'],
                chart_spec['position'],
                title=chart_spec.get('title'),
                data_selection=chart_spec.get('data_selection'),
                chart_options=chart_spec.get('options')
            )
            
        return slide
        
    def save_presentation(self, path: Union[str, Path], **kwargs) -> str:
        """
        Save the presentation.
        
        Args:
            path: Path to save the presentation
            **kwargs: Additional save options
            
        Returns:
            Actual path where the presentation was saved
        """
        path = str(path)
        self.presentation.save(path)
        return path 