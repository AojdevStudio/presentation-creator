"""
Module for finalizing presentations and managing presentation metadata.
"""
from typing import Dict, Any, Optional
from datetime import datetime
from pptx import Presentation
from pptx.presentation import Presentation as PresentationType
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

class PresentationFinalizer:
    """Handles presentation finalization and metadata management."""
    
    def __init__(self, presentation: PresentationType):
        """
        Initialize the finalizer with a presentation.
        
        Args:
            presentation: The presentation to finalize
        """
        self.presentation = presentation
        
    def add_metadata(self, metadata: Dict[str, Any]) -> None:
        """
        Add metadata to the presentation properties.
        
        Args:
            metadata: Dictionary containing metadata fields:
                     - title: Presentation title
                     - author: Author name
                     - company: Company name
                     - keywords: Keywords/tags
                     - comments: Additional comments
        """
        core_props = self.presentation.core_properties
        if 'title' in metadata:
            core_props.title = metadata['title']
        if 'author' in metadata:
            core_props.author = metadata['author']
        if 'company' in metadata:
            core_props.company = metadata['company']
        if 'keywords' in metadata:
            core_props.keywords = metadata['keywords']
        if 'comments' in metadata:
            core_props.comments = metadata['comments']
            
        # Add creation and last modified dates
        core_props.created = datetime.now()
        core_props.last_modified_by = metadata.get('author', 'Unknown')
        core_props.last_printed = datetime.now()
        core_props.modified = datetime.now()
        
    def add_title_slide(self, title: str, subtitle: Optional[str] = None) -> None:
        """
        Add or update the title slide.
        
        Args:
            title: Main title text
            subtitle: Optional subtitle text
        """
        # Get the first slide or add one if none exists
        if len(self.presentation.slides) == 0:
            slide_layout = self.presentation.slide_layouts[0]  # Title Slide layout
            slide = self.presentation.slides.add_slide(slide_layout)
        else:
            # Check if first slide is a title slide
            first_slide = self.presentation.slides[0]
            if not hasattr(first_slide.shapes, 'title'):
                # Create new title slide
                slide_layout = self.presentation.slide_layouts[0]
                new_slide = self.presentation.slides.add_slide(slide_layout)
                # Move it to the beginning
                slides = self.presentation.slides
                slides._sldIdLst.remove(new_slide._element)
                slides._sldIdLst.insert(0, new_slide._element)
                slide = new_slide
            else:
                slide = first_slide
        
        # Set title
        if slide.shapes.title:
            title_shape = slide.shapes.title
            title_shape.text = title
            # Format title
            for paragraph in title_shape.text_frame.paragraphs:
                paragraph.font.size = Pt(44)
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(33, 33, 33)
        
        # Set subtitle if provided
        if subtitle and slide.placeholders[1]:  # Second placeholder is usually subtitle
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle
            # Format subtitle
            for paragraph in subtitle_shape.text_frame.paragraphs:
                paragraph.font.size = Pt(24)
                paragraph.font.color.rgb = RGBColor(89, 89, 89)
    
    def add_footer(self, text: str) -> None:
        """
        Add footer to all slides.
        
        Args:
            text: Footer text to add
        """
        slides = list(self.presentation.slides)
        for slide in slides:
            # Skip footer on title slide (first slide)
            if slide == slides[0]:
                continue
                
            footer = slide.shapes.add_textbox(
                Inches(0.5),  # Left
                Inches(7),    # Top
                Inches(9),    # Width
                Inches(0.5)   # Height
            )
            text_frame = footer.text_frame
            
            # Clear any existing paragraphs
            for _ in range(len(text_frame.paragraphs) - 1):
                text_frame._p.remove(text_frame._p[0])
                
            # Set text in first paragraph
            paragraph = text_frame.paragraphs[0]
            paragraph.text = text
            paragraph.font.size = Pt(10)
            paragraph.font.color.rgb = RGBColor(128, 128, 128)
            
            # Configure text frame
            text_frame.word_wrap = True
            text_frame.auto_size = True
            
            # Add a custom name to identify this shape as footer
            footer.name = 'Footer Text'
    
    def add_slide_numbers(self) -> None:
        """Add slide numbers to all slides except the title slide."""
        slides = list(self.presentation.slides)  # Convert to list to avoid iterator issues
        for i, slide in enumerate(slides[1:], start=2):  # Skip first slide
            number = slide.shapes.add_textbox(
                Inches(9),    # Left
                Inches(7),    # Top
                Inches(0.5),  # Width
                Inches(0.5)   # Height
            )
            text_frame = number.text_frame
            # Clear any existing paragraphs
            for _ in range(len(text_frame.paragraphs) - 1):
                text_frame._p.remove(text_frame._p[0])
            # Set text in first paragraph
            paragraph = text_frame.paragraphs[0]
            paragraph.text = str(i)
            paragraph.font.size = Pt(10)
            paragraph.font.color.rgb = RGBColor(128, 128, 128)
            # Configure text frame
            text_frame.word_wrap = True
            text_frame.auto_size = True
    
    def finalize(self, metadata: Dict[str, Any], add_title_slide: bool = True) -> None:
        """
        Finalize the presentation by adding all necessary metadata and elements.
        
        Args:
            metadata: Dictionary containing presentation metadata
            add_title_slide: Whether to add/update the title slide
        """
        # Store existing content slides and their layouts
        content_slides = []
        for slide in self.presentation.slides:
            content_slides.append(slide.slide_layout)
            
        # Clear all slides to start fresh
        while len(self.presentation.slides) > 0:
            xml_slides = self.presentation.slides._sldIdLst
            xml_slides.remove(xml_slides[0])
            
        # Add metadata
        self.add_metadata(metadata)
        
        # Add title slide if requested
        if add_title_slide:
            self.add_title_slide(
                metadata.get('title', 'Untitled Presentation'),
                metadata.get('subtitle')
            )
            
        # Re-add content slides with their original layouts
        for layout in content_slides:
            self.presentation.slides.add_slide(layout)
            
        # Add footer if specified
        if 'footer' in metadata:
            self.add_footer(metadata['footer'])
        
        # Add slide numbers last
        self.add_slide_numbers() 